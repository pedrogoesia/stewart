"""
Geração de relatórios fotográficos em PowerPoint para a Stewart Construtora.

Usa o template oficial (template/TEMPLATE_STEWART.pptx) para herdar a identidade
visual (capa, barra vermelha, logo e layout de fotos) e monta dinamicamente:

  - 1 slide de capa com nome da obra, endereço e mês/ano do relatório;
  - N slides de fotos (2 fotos por slide) agrupados por cômodo, com a legenda
    (descrição) embaixo de cada foto e o nome do cômodo no cabeçalho.

O layout de fotos do template ("2_Em Branco") já possui dois espaços
reservados para imagem e dois para legenda. Inserimos as fotos nesses espaços
(que recortam a imagem automaticamente para preencher) e adicionamos no
cabeçalho o nome do cômodo (à esquerda) e o mês/ano (à direita).
"""

import copy

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Constantes do template (em EMU - English Metric Units), extraídas do layout
# "2_Em Branco" (slideLayout13) do template oficial.
# ---------------------------------------------------------------------------
HEADER_Y = 487481          # topo da barra vermelha do cabeçalho
HEADER_H = 400692          # altura da barra
ROOM_BOX = (Emu(415470), Emu(HEADER_Y), Emu(5000000), Emu(HEADER_H))
MONTH_BOX = (Emu(8452748), Emu(HEADER_Y), Emu(3323642), Emu(HEADER_H))

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_FONT = "Helvetica"

# Índices dos espaços reservados no layout de fotos.
PH_PIC_LEFT = 13
PH_PIC_RIGHT = 14
PH_CAP_LEFT = 15
PH_CAP_RIGHT = 16

PHOTO_LAYOUT_NAME = "2_Em Branco"
STATIC_MONTH_SHAPE = "CaixaDeTexto 8"  # mês fixo do layout, removido em runtime


def _find_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise RuntimeError(f"Layout '{name}' não encontrado no template.")


def _replace_text_keep_format(shape, new_text, font_name=None, size=None):
    """Substitui o texto de uma forma preservando a formatação do 1º run.

    Opcionalmente aplica fonte (font_name) e tamanho em pontos (size).
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = new_text
        # remove runs extras, se houver
        for extra in p.runs[1:]:
            extra._r.getparent().remove(extra._r)
        run = p.runs[0]
    else:
        run = p.add_run()
        run.text = new_text
    if font_name:
        run.font.name = font_name
    if size is not None:
        run.font.size = Pt(size)


def _set_cover(slide, obra_nome, endereco, periodo_label):
    """Preenche os textos da capa a partir dos rótulos do template."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip().upper()
        if txt.startswith("RELATÓRIO FOTOGRÁFICO"):
            _replace_text_keep_format(
                shape, f"RELATÓRIO FOTOGRÁFICO – {periodo_label}")
        elif txt == "NOME DA OBRA":
            _replace_text_keep_format(shape, obra_nome, "Helvetica", 19)
        elif txt == "ENDEREÇO DA OBRA":
            _replace_text_keep_format(shape, endereco or "", "Helvetica", 9)


def _remove_static_month(prs):
    """Remove o texto de mês fixo do layout de fotos para inserirmos o nosso."""
    layout = _find_layout(prs, PHOTO_LAYOUT_NAME)
    for shape in list(layout.shapes):
        if shape.name == STATIC_MONTH_SHAPE:
            shape._element.getparent().remove(shape._element)


def _add_header_label(slide, box, text, align):
    left, top, width, height = box
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.name = HEADER_FONT
    run.font.color.rgb = WHITE
    return tb


def _placeholder(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def _delete_slide(prs, index):
    """Remove um slide por completo (entrada da lista, relação e part)."""
    sldId = prs.slides._sldIdLst[index]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sldId)


def _set_caption(slide, idx, text):
    """
    Define a legenda da foto com formatação:
      - tudo em negrito;
      - o nome do cômodo (parte antes do primeiro " - ") também sublinhado;
      - primeira letra em maiúscula.
    """
    ph = _placeholder(slide, idx)
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]

    text = (text or "").strip()
    if not text:
        return
    text = text[0].upper() + text[1:]

    # O nome do cômodo é a parte antes do primeiro "-". Ele fica sempre em
    # negrito + sublinhado; o restante (incluindo o "-") fica em negrito.
    pos = text.find("-")
    if pos != -1:
        nome = text[:pos].rstrip()
        resto = text[pos:]            # inclui o "-" e o que vem depois
        r1 = p.add_run()
        r1.text = nome
        r1.font.bold = True
        r1.font.underline = True
        r2 = p.add_run()
        r2.text = " " + resto
        r2.font.bold = True
    else:
        r = p.add_run()
        r.text = text
        r.font.bold = True
        r.font.underline = True


def _add_photo_slide(prs, layout, comodo_nome, periodo_label, fotos_par):
    """Cria um slide com até 2 fotos (lista de dicts {path, descricao})."""
    slide = prs.slides.add_slide(layout)

    # Cabeçalho: apenas mês/ano à direita (o cômodo já aparece nas legendas).
    _add_header_label(slide, MONTH_BOX, periodo_label, PP_ALIGN.RIGHT)

    # Foto da esquerda (sempre existe)
    f1 = fotos_par[0]
    pic_l = _placeholder(slide, PH_PIC_LEFT)
    pic_l.insert_picture(f1["path"])
    _set_caption(slide, PH_CAP_LEFT, f1.get("descricao", ""))

    # Foto da direita (opcional)
    if len(fotos_par) > 1:
        f2 = fotos_par[1]
        pic_r = _placeholder(slide, PH_PIC_RIGHT)
        pic_r.insert_picture(f2["path"])
        _set_caption(slide, PH_CAP_RIGHT, f2.get("descricao", ""))
    else:
        # Remove os espaços reservados não usados para um slide limpo.
        for idx in (PH_PIC_RIGHT, PH_CAP_RIGHT):
            ph = _placeholder(slide, idx)
            if ph is not None:
                _remove_shape(ph)

    return slide


def gerar_relatorio(template_path, output_path, obra, periodo_label, comodos):
    """
    Gera o relatório .pptx.

    Parâmetros
    ----------
    template_path : str   caminho do template .pptx
    output_path   : str   caminho de saída do relatório
    obra          : dict  {"nome": str, "endereco": str}
    periodo_label : str   ex.: "JUNHO 2026"
    comodos       : list  [{"nome": str,
                            "fotos": [{"path": str, "descricao": str}, ...]}]
    """
    prs = Presentation(template_path)

    _remove_static_month(prs)
    layout = _find_layout(prs, PHOTO_LAYOUT_NAME)

    # slide[0] = capa (reutilizada); slide[1] = exemplo (removido)
    _set_cover(prs.slides[0], obra["nome"], obra.get("endereco", ""),
               periodo_label)
    while len(prs.slides._sldIdLst) > 1:
        _delete_slide(prs, 1)

    for comodo in comodos:
        fotos = comodo.get("fotos", [])
        if not fotos:
            continue
        for i in range(0, len(fotos), 2):
            _add_photo_slide(prs, layout, comodo["nome"], periodo_label,
                             fotos[i:i + 2])

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    # Teste rápido com as imagens de exemplo do template.
    import os
    here = os.path.dirname(os.path.abspath(__file__))
    tpl = os.path.join(here, "template", "TEMPLATE_STEWART.pptx")
    out = "/tmp/relatorio_teste.pptx"
    gerar_relatorio(
        tpl, out,
        obra={"nome": "RESIDENCIAL TESTE",
              "endereco": "Rua das Obras, 123 - São Paulo/SP"},
        periodo_label="JUNHO 2026",
        comodos=[
            {"nome": "SALA DE ESTAR", "fotos": [
                {"path": "/tmp/sample1.jpg", "descricao": "Parede norte - reboco concluído"},
                {"path": "/tmp/sample2.jpg", "descricao": "Vista geral do ambiente"},
                {"path": "/tmp/sample1.jpg", "descricao": "Foto ímpar - sozinha no slide"},
            ]},
            {"nome": "COZINHA", "fotos": [
                {"path": "/tmp/sample2.jpg", "descricao": "Bancada instalada"},
                {"path": "/tmp/sample1.jpg", "descricao": ""},
            ]},
        ],
    )
    print("Gerado:", out, os.path.getsize(out), "bytes")
